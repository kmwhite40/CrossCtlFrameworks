"""Slide deck parser.

Slide number is the citation an assessor can act on, and the title placeholder is
the only reliable section label a deck offers, so it becomes both the page's
``section_title`` and the heading path for every block on that slide. Shapes are
walked in the order python-pptx exposes them; the title shape is emitted first so
a slide reads title-then-body.
"""

from __future__ import annotations

import io

from pptx import Presentation

from ...logging import get_logger
from .base import ParsedBlock, ParsedDocument, ParsedPage

log = get_logger(__name__)

PARSER_NAME = "pptx"
MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def parse_pptx(data: bytes, filename: str, media_type: str = MEDIA_TYPE) -> ParsedDocument:
    """Parse a .pptx into one page per slide of title-scoped text blocks."""
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:  # a malformed deck must not kill the run
        log.warning("prep.parse.pptx_failed", filename=filename, error=str(exc))
        return ParsedDocument(
            filename=filename, media_type=media_type, parser_name=PARSER_NAME,
            error=f"could not open presentation: {exc}",
        )

    pages: list[ParsedPage] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = (title_shape.text or "").strip() if title_shape is not None else ""
        heading_path = [title] if title else []

        ordered = (
            [title_shape, *[s for s in slide.shapes if s is not title_shape]]
            if title_shape is not None
            else list(slide.shapes)
        )
        blocks: list[ParsedBlock] = []
        for shape_index, shape in enumerate(ordered):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            blocks.append(
                ParsedBlock(
                    block_id=f"s{slide_number}shape{shape_index}",
                    block_type="slide_text",
                    text=text,
                    heading_path=list(heading_path),
                )
            )
        pages.append(
            ParsedPage(
                page_number=slide_number, blocks=blocks, section_title=title or None
            )
        )

    return ParsedDocument(
        filename=filename, media_type=media_type, parser_name=PARSER_NAME, pages=pages
    )
