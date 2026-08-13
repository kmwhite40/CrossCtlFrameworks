"""PPTX parser — slide numbering and title-derived sections."""

from __future__ import annotations

import io
import logging

from pptx import Presentation
from pptx.util import Inches

from ccf.prep.parsers import dispatch
from ccf.prep.parsers.pptx import parse_pptx


def _pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Access Control Architecture"
    slide.placeholders[1].text = "All administrative access requires MFA."
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Audit Pipeline"
    box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "Logs ship to the SIEM within five minutes."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_body_text_carries_its_slide_number() -> None:
    doc = parse_pptx(_pptx_bytes(), "brief.pptx")
    line = next(x for x in doc.iter_lines() if "SIEM" in x.content)
    assert line.page_number == 2


def test_slide_title_becomes_the_section_path() -> None:
    doc = parse_pptx(_pptx_bytes(), "brief.pptx")
    line = next(x for x in doc.iter_lines() if "MFA" in x.content)
    assert line.section_path == "Access Control Architecture"


def test_title_is_not_emitted_twice() -> None:
    """Regression: slide.shapes.title is a fresh wrapper object each call, distinct

    by identity from the object yielded while iterating slide.shapes even though
    both wrap the same XML element. Filtering "the remaining shapes" by object
    identity (`is not title_shape`) therefore never excludes the real title, so it
    lands in the shape order twice — once explicitly first, once again from the
    iteration — and every titled slide emits a duplicate title line. This asserts
    the exact ordered content list for slide 1, length included, so a
    reintroduced duplicate fails instead of being masked by a substring check.
    """
    doc = parse_pptx(_pptx_bytes(), "brief.pptx")
    slide_one = [x.content for x in doc.iter_lines() if x.page_number == 1]
    assert slide_one == [
        "Access Control Architecture",
        "All administrative access requires MFA.",
    ]


def test_section_title_is_recorded_on_the_page() -> None:
    doc = parse_pptx(_pptx_bytes(), "brief.pptx")
    assert [p.section_title for p in doc.pages] == [
        "Access Control Architecture",
        "Audit Pipeline",
    ]


def test_block_type_marks_slide_text() -> None:
    doc = parse_pptx(_pptx_bytes(), "brief.pptx")
    assert all(x.block_type == "slide_text" for x in doc.iter_lines())


def test_corrupt_pptx_returns_an_error_document_rather_than_raising() -> None:
    doc = parse_pptx(b"not a pptx", "broken.pptx")
    assert not doc.success
    assert doc.error is not None


def test_dispatch_routes_pptx() -> None:
    doc = dispatch(_pptx_bytes(), "brief.pptx", None)
    assert doc.parser_name == "pptx"


def test_corrupt_pptx_failure_path_survives_an_enabled_logger() -> None:
    """Regression for the reserved-LogRecord-attribute defect.

    tests/conftest.py runs Alembic migrations, and migrations/env.py calls
    logging.config.fileConfig(..., disable_existing_loggers=True), which disables
    this module's stdlib logger for the whole pytest session. A logger.warning()
    call on a *disabled* logger short-circuits before building a LogRecord, so a
    bug in the log call itself (e.g. passing `extra={"filename": ...}`, which
    collides with LogRecord's own reserved `filename` attribute and raises
    KeyError from stdlib logging) would never fire and this suite would pass
    while the same call raises outside pytest. Force the logger enabled so this
    test genuinely exercises the log call instead of being masked by that side
    effect.
    """
    logger = logging.getLogger("ccf.prep.parsers.pptx")
    original = logger.disabled
    logger.disabled = False
    try:
        doc = parse_pptx(b"not a pptx", "broken.pptx")
    finally:
        logger.disabled = original
    assert not doc.success
    assert doc.error is not None
